"""
蓝心校园AI管家 - 答辩PPT蓝白主题转换脚本
将深色主题的PPT 2转换为蓝白配色，并替换团队页占位符
"""
import copy
import shutil
from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
import io, sys, os

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

SRC = '蓝心校园AI管家-比赛答辩展示.pptx'
DST = '蓝心校园AI管家-最终答辩.pptx'

# ========== 颜色映射表 ==========
# 将深色主题的 sRGB 颜色映射到蓝白主题
COLOR_MAP = {
    # 深色背景 → 白色/浅蓝
    '07111F': 'FFFFFF',  # 最深背景 → 纯白
    '0D2C5A': 'F0F5FF',  # 深蓝背景 → 极浅蓝
    '0B243A': 'F8FAFE',  # 深蓝黑 → 近白
    '0B2038': 'F8FAFE',  # 深蓝黑变体 → 近白
    '071827': 'FFFFFF',  # 深色 → 纯白
    '091B30': 'F0F5FF',  # 深色 → 浅蓝
    '0C2037': 'EFF6FF',  # 卡片深色背景 → 浅蓝卡片
    '14345D': 'DBEAFE',  # 深蓝 → 中浅蓝
    '1D496A': 'BFDBFE',  # 深蓝 → 浅蓝边框
    '163F3C': 'EFF6FF',  # 深绿 → 浅蓝
    '23658A': 'DBEAFE',  # 深青 → 浅蓝

    # 青色强调 → 深蓝主色
    '28D7FF': '1E40AF',  # 亮青 → 深蓝(主色)

    # 绿色强调 → 中蓝
    '42E6A4': '2563EB',  # 亮绿 → 亮蓝(辅助色)

    # 浅色文字(深底上) → 深色文字(白底上)
    'EEF7FF': '1E293B',  # 近白文字 → 深灰文字

    # 灰色文字 → 中灰
    '9DB6CF': '475569',  # 浅灰蓝 → 深灰蓝
    '5F7895': '64748B',  # 中灰蓝 → 石板灰

    # 粉色/红色强调 → 蓝色
    'FF5C7A': '3B82F6',  # 粉红 → 中蓝

    # 橙色强调 → 浅蓝
    'FFB65C': '60A5FA',  # 橙色 → 浅蓝
}

def remap_color(hex_val):
    """将颜色值映射到蓝白主题，保留不在映射表中的颜色"""
    upper = hex_val.upper()
    return COLOR_MAP.get(upper, upper)

def process_element(el):
    """递归处理XML元素中的颜色属性"""
    for child in el.iter():
        # 处理 srgbClr
        for color_el in child.findall('.//{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr'):
            val = color_el.get('val')
            if val:
                new_val = remap_color(val)
                if new_val != val:
                    color_el.set('val', new_val)

        # 处理 schemeClr 的 lastClr
        for color_el in child.findall('.//{http://schemas.openxmlformats.org/drawingml/2006/main}schemeClr'):
            last_clr = color_el.get('lastClr')
            if last_clr:
                new_val = remap_color(last_clr)
                if new_val != last_clr:
                    color_el.set('lastClr', new_val)

def convert_shape_colors(shape):
    """转换单个形状中的所有颜色"""
    process_element(shape._element)

# ========== 步骤1: 复制源文件 ==========
print(f"复制 {SRC} → {DST}")
shutil.copy2(SRC, DST)

# ========== 步骤2: 打开并处理 ==========
print("打开演示文稿...")
prs = Presentation(DST)

# ========== 步骤3: 修改主题 XML ==========
print("修改主题配色方案...")
# 通过操作压缩包内的 theme XML 来修改主题
import zipfile
import tempfile
import os

def modify_theme_xml(pptx_path):
    """直接修改 PPTX 压缩包内的主题 XML"""
    tmp_path = pptx_path + '.tmp'

    with zipfile.ZipFile(pptx_path, 'r') as zin:
        with zipfile.ZipFile(tmp_path, 'w', zipfile.ZIP_DEFLATED) as zout:
            for item in zin.infolist():
                data = zin.read(item.filename)

                if 'theme' in item.filename.lower() and item.filename.endswith('.xml'):
                    # 修改主题配色
                    root = etree.fromstring(data)
                    nsmap = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}

                    # 新配色方案：蓝白主题
                    new_scheme = {
                        'dk1': '000000',     # 黑色(用于深色文字)
                        'lt1': 'FFFFFF',     # 白色(用于背景)
                        'dk2': '1E3A5F',     # 深蓝
                        'lt2': 'DBEAFE',     # 浅蓝
                        'accent1': '1E40AF', # 深蓝(主色)
                        'accent2': '2563EB', # 亮蓝
                        'accent3': '3B82F6', # 中蓝
                        'accent4': '60A5FA', # 浅蓝
                        'accent5': '93C5FD', # 更浅蓝
                        'accent6': 'DBEAFE', # 最浅蓝
                        'hlink': '2563EB',   # 链接蓝
                        'folHlink': '1E3A5F',# 访问后链接
                    }

                    for cs in root.iter('{http://schemas.openxmlformats.org/drawingml/2006/main}clrScheme'):
                        cs.set('name', 'BlueWhite')
                        for child in cs:
                            tag = child.tag.split('}')[-1]
                            if tag in new_scheme:
                                for sub in child:
                                    subtag = sub.tag.split('}')[-1]
                                    if subtag == 'srgbClr':
                                        sub.set('val', new_scheme[tag])
                                    elif subtag == 'sysClr':
                                        # 更新 lastClr
                                        sub.set('lastClr', new_scheme[tag])

                    data = etree.tostring(root, xml_declaration=True, encoding='UTF-8', standalone=True)

                zout.writestr(item, data)

    # 替换原文件
    os.replace(tmp_path, pptx_path)

modify_theme_xml(DST)

# 重新打开（因为zip修改了）
prs = Presentation(DST)

# ========== 步骤4: 转换所有形状颜色 ==========
print("转换形状颜色...")
slides = list(prs.slides)
total_shapes = 0
for i, slide in enumerate(slides):
    for shape in slide.shapes:
        convert_shape_colors(shape)
        total_shapes += 1
    print(f"  处理幻灯片 {i+1}/{len(slides)}")

print(f"共处理 {total_shapes} 个形状")

# ========== 步骤5: 修改幻灯片背景 ==========
print("修改幻灯片背景...")
for i, slide in enumerate(slides):
    # 设置背景为纯白
    bg = slide.background
    # 通过 XML 设置背景
    cSld = slide._element
    # 查找或创建 bg 元素
    nsmap = {
        'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
        'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    }

    # 移除现有背景
    for existing_bg in cSld.findall('{http://schemas.openxmlformats.org/presentationml/2006/main}cSld'):
        # cSld 本身不需要修改
        pass

    # 通过 slide element 找 bg
    bg_elements = slide._element.findall('{http://schemas.openxmlformats.org/presentationml/2006/main}cSld/{http://schemas.openxmlformats.org/presentationml/2006/main}bg')
    for bg_el in bg_elements:
        slide._element.find('{http://schemas.openxmlformats.org/presentationml/2006/main}cSld').remove(bg_el)

# ========== 步骤6: 替换团队页(Slide 2) ==========
print("替换团队页内容...")
team_slide = slides[1]  # Slide index 1 = Slide 2

# 团队成员的姓名和角色（从PPT 1封面和团队页提取）
# 用户可根据实际情况修改
team_members = [
    {
        'name': '产品经理',
        'role': '统筹与规划',
        'desc': '负责模块：项目整体规划与资源协调 | 贡献亮点：产品方向把控、版本发布管理 | 答辩职责：产品介绍与愿景阐述'
    },
    {
        'name': '后端开发',
        'role': '架构与API',
        'desc': '负责模块：Spring Boot架构搭建、核心API开发 | 贡献亮点：SSE流式输出、JWT认证体系 | 答辩职责：技术架构讲解'
    },
    {
        'name': '前端开发',
        'role': '交互与界面',
        'desc': '负责模块：原生前端界面实现、交互体验打磨 | 贡献亮点：Markdown渲染、OCR识别集成 | 答辩职责：功能演示操作'
    },
    {
        'name': 'AI算法工程师',
        'role': '模型与RAG',
        'desc': '负责模块：RAG检索增强、大模型调用调优 | 贡献亮点：vivo AIGC接口对接、知识库构建 | 答辩职责：AI能力讲解'
    },
    {
        'name': 'UI/UX设计师',
        'role': '视觉与体验',
        'desc': '负责模块：全链路用户视觉设计、交互逻辑 | 贡献亮点：品牌视觉规范、用户体验研究 | 答辩职责：设计理念陈述'
    },
]

# 在 Slide 2 中找到对应文本形状并替换
# 成员姓名: shapes 11, 16, 21, 26, 31
# 角色分工: shapes 12, 17, 22, 27, 32
# 描述文字: shapes 13, 18, 23, 28, 33
name_indices = [11, 16, 21, 26, 31]
role_indices = [12, 17, 22, 27, 32]
desc_indices = [13, 18, 23, 28, 33]

shapes_list = list(team_slide.shapes)

def clear_and_set_text(shape, text, color_rgb):
    """彻底清除形状中的所有段落并设置新文本"""
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    # 清除所有段落
    for para in tf.paragraphs:
        para.clear()
    # 确保至少有一个段落
    if len(tf.paragraphs) == 0:
        tf.add_paragraph()
    run = tf.paragraphs[0].add_run()
    run.text = text
    run.font.color.rgb = color_rgb

for idx, member in enumerate(team_members):
    if idx < len(name_indices):
        ni = name_indices[idx]
        if ni < len(shapes_list):
            clear_and_set_text(shapes_list[ni], member['name'], RGBColor(0x1E, 0x29, 0x3B))

    if idx < len(role_indices):
        ri = role_indices[idx]
        if ri < len(shapes_list):
            clear_and_set_text(shapes_list[ri], member['role'], RGBColor(0x1E, 0x40, 0xAF))

    if idx < len(desc_indices):
        di = desc_indices[idx]
        if di < len(shapes_list):
            clear_and_set_text(shapes_list[di], member['desc'], RGBColor(0x47, 0x55, 0x69))

# 修改提示文字
if 34 < len(shapes_list):
    clear_and_set_text(shapes_list[34],
        '请将姓名替换为实际成员姓名，角色分工可根据实际情况调整',
        RGBColor(0x64, 0x74, 0x8B))

# ========== 步骤7: 修改底部页脚文字颜色 ==========
print("调整页脚样式...")
for slide in slides:
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text
            if '蓝心校园 AI 管家' in text and len(text) < 30:
                # 这是页脚文字，改颜色
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        run.font.color.rgb = RGBColor(0x64, 0x74, 0x8B)  # 灰蓝

# ========== 步骤8: 保存 ==========
print(f"保存到 {DST}...")
prs.save(DST)
print(f"✅ 完成！最终答辩PPT已保存到: {DST}")
print(f"   幻灯片数量: {len(slides)}")
