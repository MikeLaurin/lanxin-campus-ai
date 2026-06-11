"""
蓝心校园AI管家 - 答辩PPT生成器
从PPT 1提取内容，压缩为13页适合5分钟答辩的精简版
保留原Office主题配色，优化排版
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import io, sys

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

# ========== 配色方案 (保留原PPT 1的Office主题色) ==========
BLUE      = RGBColor(0x44, 0x72, 0xC4)  # accent1 主色
DARK_BLUE = RGBColor(0x2B, 0x50, 0x8A)  # 深蓝 标题
ORANGE    = RGBColor(0xED, 0x7D, 0x31)  # accent2 强调
LIGHT_BLUE= RGBColor(0x5B, 0x9B, 0xD5)  # accent5 浅蓝
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
BLACK     = RGBColor(0x00, 0x00, 0x00)
DARK_GRAY = RGBColor(0x44, 0x54, 0x6A)  # dk2
LIGHT_GRAY= RGBColor(0xE7, 0xE6, 0xE6)  # lt2
MID_GRAY  = RGBColor(0x6B, 0x72, 0x80)
BG_WHITE  = RGBColor(0xFA, 0xFA, 0xFA)

# 幻灯片尺寸 (16:9)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# ========== 工具函数 ==========
def add_blank_slide():
    """添加空白幻灯片"""
    layout = prs.slide_layouts[6]  # blank layout
    return prs.slides.add_slide(layout)

def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=DARK_GRAY, bold=False, alignment=PP_ALIGN.LEFT,
                font_name='Microsoft YaHei'):
    """添加文本框"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf

def add_shape_bg(slide, left, top, width, height, color, shape_type=MSO_SHAPE.RECTANGLE):
    """添加纯色矩形背景"""
    shape = slide.shapes.add_shape(shape_type, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()  # no border
    return shape

def add_accent_bar(slide, left, top, width, height, color=BLUE):
    """添加强调色条"""
    return add_shape_bg(slide, left, top, width, height, color)

def add_page_number(slide, num, total=13):
    """右下角页码"""
    add_textbox(slide, 12.0, 7.0, 1.2, 0.4, f'{num}/{total}',
                font_size=10, color=MID_GRAY, alignment=PP_ALIGN.RIGHT)

def add_slide_title(slide, title, subtitle=None):
    """统一标题栏：左侧蓝色竖条 + 标题"""
    # 左侧蓝色竖条
    add_accent_bar(slide, 0.6, 0.5, 0.08, 0.6, BLUE)
    # 标题
    add_textbox(slide, 1.0, 0.45, 10, 0.7, title, font_size=32,
                color=DARK_BLUE, bold=True)
    if subtitle:
        add_textbox(slide, 1.0, 1.1, 10, 0.5, subtitle, font_size=14,
                    color=MID_GRAY)
    # 底部分隔线
    add_accent_bar(slide, 1.0, 1.55, 11.3, 0.015, LIGHT_GRAY)

def add_card(slide, left, top, width, height, title, body, icon_color=BLUE):
    """添加卡片：圆角矩形 + 标题 + 正文"""
    # 背景
    card = add_shape_bg(slide, left, top, width, height, BG_WHITE, MSO_SHAPE.ROUNDED_RECTANGLE)
    # 顶部色条
    add_accent_bar(slide, left, top, width, 0.06, icon_color)
    # 标题
    add_textbox(slide, left + 0.3, top + 0.25, width - 0.6, 0.45, title,
                font_size=16, color=DARK_BLUE, bold=True)
    # 正文
    add_textbox(slide, left + 0.3, top + 0.7, width - 0.6, height - 0.9, body,
                font_size=12, color=MID_GRAY)

# ========== 幻灯片 1: 封面 ==========
print("创建 Slide 1: 封面")
slide = add_blank_slide()
# 背景色块
add_shape_bg(slide, 0, 0, 13.333, 7.5, DARK_BLUE)

# 装饰线条
add_accent_bar(slide, 1.5, 2.0, 10.333, 0.03, LIGHT_BLUE)
add_accent_bar(slide, 1.5, 5.2, 10.333, 0.03, LIGHT_BLUE)

# 主标题
add_textbox(slide, 1.5, 2.3, 10.333, 1.4, '蓝心校园 AI 管家——课伴Pro',
            font_size=44, color=WHITE, bold=True)

# 副标题
add_textbox(slide, 1.5, 3.7, 10.333, 0.8,
            '你的校园学习搭子，更懂你的 AI 管家\n用智能技术重塑高效学习新体验',
            font_size=20, color=LIGHT_BLUE)

# 团队和日期
add_textbox(slide, 1.5, 5.5, 10.333, 0.5,
            '汇报人：剑研风华团队（林家圳  温加研  张剑  马昌赫  程宇啸）',
            font_size=16, color=WHITE)
add_textbox(slide, 1.5, 6.0, 10.333, 0.4,
            '2026年06月08日',
            font_size=14, color=LIGHT_BLUE)

# ========== 幻灯片 2: 项目简介（压缩原PPT 4+5+7+8页） ==========
print("创建 Slide 2: 项目简介")
slide = add_blank_slide()
add_slide_title(slide, '项目简介', 'Project Introduction — 探索项目初衷，明确核心目标与价值愿景')
add_page_number(slide, 2)

# 定义 + 一句话核心
add_textbox(slide, 1.0, 1.8, 7.5, 0.8,
            '新一代大学生AI学习伴侣',
            font_size=24, color=DARK_BLUE, bold=True)
add_textbox(slide, 1.0, 2.4, 7.5, 0.6,
            '以课堂端侧多模态智能笔记为核心，结合"朋友+管家"式时间管理，\n构建"信息捕捉→知识沉淀→行动复盘"的全链路AI驱动学习闭环。',
            font_size=14, color=DARK_GRAY)

# 右侧：产品定位三要素
add_shape_bg(slide, 9.0, 1.9, 3.8, 1.4, BG_WHITE, MSO_SHAPE.ROUNDED_RECTANGLE)
add_textbox(slide, 9.2, 2.0, 3.4, 0.3, '产品形态', font_size=11, color=MID_GRAY, bold=True)
add_textbox(slide, 9.2, 2.25, 3.4, 0.25, 'VIVO快应用', font_size=14, color=DARK_BLUE, bold=True)
add_textbox(slide, 9.2, 2.48, 3.4, 0.6, '免安装、即点即用\n轻量化触达大学生群体', font_size=10, color=MID_GRAY)
add_textbox(slide, 9.2, 2.95, 3.4, 0.25, '品牌人设：懂你的"校园朋友"小蓝', font_size=11, color=ORANGE, bold=True)

# 4个核心能力（紧凑卡片）
cards = [
    ('课堂多模态智能笔记', '拍照、录音一键生成结构化笔记\n自动提炼重点知识与核心脉络'),
    ('时间管理大师', '自动解析教务通知、作业DDL\n课程预警及复习节点提醒'),
    ('朋友式AI交互', '温暖活泼的"朋友式"多模态对话\n鼓励式反馈，缓解学习焦虑'),
    ('任务自动化处理', '自动归档笔记、生成复习资料\n学习周报，省心更高效'),
]
for i, (title, body) in enumerate(cards):
    x = 1.0 + i * 3.0
    add_card(slide, x, 3.6, 2.8, 2.2, title, body, [BLUE, ORANGE, LIGHT_BLUE, DARK_BLUE][i])

# 底部：愿景
add_accent_bar(slide, 1.0, 6.15, 11.3, 0.015, LIGHT_GRAY)
add_textbox(slide, 1.0, 6.35, 11.3, 0.5,
            '产品愿景：致力于成为每个大学生的专属AI学习伙伴，提供高效、温暖且高度落地的智慧学习解决方案，重塑校园学习体验。',
            font_size=12, color=ORANGE, bold=True, alignment=PP_ALIGN.CENTER)

# ========== 幻灯片 3: 研发背景与目标人群（压缩原PPT 7+8页） ==========
print("创建 Slide 3: 研发背景与目标人群")
slide = add_blank_slide()
add_slide_title(slide, '研发背景与目标人群', '聚焦学习效率待提升的在校大学生，提供针对性解决方案')
add_page_number(slide, 3)

# 上半部分：背景四要素
add_textbox(slide, 1.0, 1.8, 11.3, 0.4, '📊 研发背景', font_size=18, color=DARK_BLUE, bold=True)

bg_items = [
    ('市场数据', '《2024中国大学生学习行为报告》显示：\n大学生普遍面临信息过载与效率低下'),
    ('国家战略', '教育部明确推动AI与教育教学\n深度融合，提供政策支持'),
    ('端侧大模型趋势', 'AIGC技术日趋成熟，端侧部署\n让模型调用更即时、更私密'),
    ('VIVO生态优势', '覆盖超7亿终端设备，免安装即点即用\n极大降低用户使用门槛'),
]
for i, (title, body) in enumerate(bg_items):
    x = 1.0 + i * 3.0
    add_card(slide, x, 2.3, 2.8, 2.0, title, body, BLUE)

# 下半部分：目标人群 + 产品边界
add_accent_bar(slide, 1.0, 4.6, 11.3, 0.015, LIGHT_GRAY)
add_textbox(slide, 1.0, 4.75, 6.0, 0.4, '🎯 目标人群', font_size=18, color=DARK_BLUE, bold=True)

targets = [
    '核心年龄层：18-24岁大一至大四本科生',
    '核心场景：课堂听讲、日常自习、期末考前冲刺',
    '功能诉求：提效减负，"不挂科、拿高分"',
    '情感诉求：在枯燥学习中，渴望陪伴感与正向鼓励',
]
for i, t in enumerate(targets):
    add_textbox(slide, 1.0, 5.2 + i * 0.35, 7.0, 0.3, f'• {t}', font_size=12, color=DARK_GRAY)

# 产品边界（右侧）
add_shape_bg(slide, 8.5, 4.75, 4.3, 2.0, BG_WHITE, MSO_SHAPE.ROUNDED_RECTANGLE)
add_accent_bar(slide, 8.5, 4.75, 4.3, 0.06, ORANGE)
add_textbox(slide, 8.8, 4.95, 3.8, 0.3, '明确边界 · 拒绝同质化', font_size=14, color=ORANGE, bold=True)
add_textbox(slide, 8.8, 5.35, 3.8, 1.3,
            '主动舍弃食堂菜单、快递查询、\n天气查询等非核心冗余功能，\n不做大而全的"校园工具箱"，\n专注深耕学习提效与陪伴的核心价值。',
            font_size=11, color=DARK_GRAY)

# ========== 幻灯片 4: 团队介绍 ==========
print("创建 Slide 4: 团队介绍")
slide = add_blank_slide()
add_slide_title(slide, '项目团队', '剑研风华 — 产品·工程·算法·设计协同')
add_page_number(slide, 4)

team = [
    ('产品经理', '统筹项目规划与需求分析\n主导产品设计与全周期管理', BLUE),
    ('后端开发', '服务器架构搭建与API开发\n数据库设计、系统底层逻辑', DARK_BLUE),
    ('前端开发', '应用界面精细化实现\n交互细节打磨、性能优化', BLUE),
    ('AI算法工程师', 'RAG模型训练与调优\nAI功能系统集成与创新', ORANGE),
    ('UI/UX设计师', '全链路视觉与交互设计\n品牌规范与用户体验研究', LIGHT_BLUE),
]
for i, (role, desc, color) in enumerate(team):
    x = 0.8 + i * 2.4
    # 头像占位圆
    add_shape_bg(slide, x + 0.6, 2.2, 1.2, 1.2, color, MSO_SHAPE.OVAL)
    add_textbox(slide, x + 0.6, 2.5, 1.2, 0.6, role[:2], font_size=22,
                color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    # 角色名
    add_textbox(slide, x, 3.7, 2.4, 0.4, role, font_size=15,
                color=DARK_BLUE, bold=True, alignment=PP_ALIGN.CENTER)
    # 描述
    add_textbox(slide, x, 4.2, 2.4, 1.0, desc, font_size=11,
                color=MID_GRAY, alignment=PP_ALIGN.CENTER)

# 底部总结
add_textbox(slide, 1.0, 5.8, 11.3, 0.5,
            '五位核心成员各展所长、紧密协作，以专业能力为基石，共同打造高效、智能、有温度的校园AI管家产品体系。',
            font_size=13, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

# ========== 幻灯片 5: 技术架构 ==========
print("创建 Slide 5: 技术架构")
slide = add_blank_slide()
add_slide_title(slide, '技术架构总览', 'Spring Boot + 原生前端 + 蓝心AIGC')
add_page_number(slide, 5)

# 分层架构图
layers = [
    ('前端交互层', '原生 HTML/CSS/JS\nMarkdown/KaTeX渲染\n语音识别/录音', BLUE),
    ('Web接入层', '请求路由 & 身份校验\n流量限制 & 参数验证\n统一异常处理', DARK_BLUE),
    ('业务服务层', 'AuthService\nLanxinApiClient\nRagService', LIGHT_BLUE),
    ('数据持久层', 'JPA + H2 数据库\n实体映射 & 查询优化\n数据一致性保障', BLUE),
    ('AI能力层', '蓝心大模型 API\nRAG检索增强\n多模态识别', ORANGE),
]
for i, (name, desc, color) in enumerate(layers):
    y = 2.0 + i * 1.0
    # 层名标签
    add_shape_bg(slide, 1.0, y, 2.3, 0.8, color, MSO_SHAPE.ROUNDED_RECTANGLE)
    add_textbox(slide, 1.0, y + 0.15, 2.3, 0.5, name, font_size=14,
                color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    # 描述
    add_textbox(slide, 3.6, y + 0.05, 8.5, 0.7, desc, font_size=12, color=DARK_GRAY)
    # 连接线
    if i < len(layers) - 1:
        pass  # 用箭头或留白表示

# 右侧关键技术点
add_textbox(slide, 1.0, 7.0, 11.3, 0.4,
            '🔑 关键技术：JWT双令牌认证 | SSE流式输出 | RAG向量检索 | 端侧优先·云端异步 | 统一异常处理',
            font_size=12, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

# ========== 幻灯片 6: 智能笔记 ==========
print("创建 Slide 6: 核心功能-智能笔记")
slide = add_blank_slide()
add_slide_title(slide, '核心功能 ① 多模态智能笔记', '课堂端侧AI，零散信息→结构化知识图谱')
add_page_number(slide, 6)

# 三种输入方式
add_textbox(slide, 1.0, 2.0, 11.3, 0.5, '📥 三种创建方式', font_size=18, color=DARK_BLUE, bold=True)
inputs = [
    ('✍️ 手写笔记', '支持Markdown实时预览\n数学公式KaTeX渲染', BLUE),
    ('📸 拍照识别', 'OCR精准提取文字与公式\n课堂板书→可编辑数字内容', ORANGE),
    ('📄 文档上传', '支持PDF/DOCX/TXT解析\n自动提取并结构化', LIGHT_BLUE),
]
for i, (title, body, color) in enumerate(inputs):
    add_card(slide, 1.0 + i * 3.9, 2.6, 3.6, 2.2, title, body, color)

# AI能力
add_textbox(slide, 1.0, 5.2, 11.3, 0.5, '🤖 AI结构化能力', font_size=18, color=DARK_BLUE, bold=True)
ai_features = '自动生成摘要 · 提取关键知识点与公式 · 智能标签分类 · 思维导图生成 · 自动写入RAG知识库'
add_textbox(slide, 1.0, 5.7, 11.3, 0.6, ai_features, font_size=14, color=DARK_GRAY)

add_textbox(slide, 1.0, 6.5, 11.3, 0.5,
            '💡 核心价值：将课堂零散信息自动转化为结构化知识，让笔记不再"记了等于没记"',
            font_size=14, color=ORANGE, bold=True, alignment=PP_ALIGN.CENTER)

# ========== 幻灯片 7: DDL时间管理 ==========
print("创建 Slide 7: 核心功能-DDL时间管理")
slide = add_blank_slide()
add_slide_title(slide, '核心功能 ② 时间管理大师', '朋友式督促，让DDL不再被遗忘')
add_page_number(slide, 7)

# 功能特点
features = [
    ('🧠 AI智能解析', '自然语言批量输入\nAI自动提取标题、日期、优先级\n用户确认后保存，降低误判风险'),
    ('📊 智能分类预警', '自动分类：考试/作业/体测/活动/论文\n首页醒目色块预警过期任务\n底部红点标记紧急程度'),
    ('🔔 主动提醒', '到期自动推送通知\n课程、考试、复习节点全覆盖\n朋友式语气，温暖不冰冷'),
]
for i, (title, body) in enumerate(features):
    y = 2.0 + i * 1.7
    add_shape_bg(slide, 1.0, y, 11.3, 1.4, BG_WHITE, MSO_SHAPE.ROUNDED_RECTANGLE)
    add_accent_bar(slide, 1.0, y, 0.08, 1.4, [BLUE, ORANGE, LIGHT_BLUE][i])
    add_textbox(slide, 1.5, y + 0.15, 3.0, 0.4, title, font_size=18, color=DARK_BLUE, bold=True)
    add_textbox(slide, 1.5, y + 0.65, 10.5, 0.7, body, font_size=14, color=DARK_GRAY)

# ========== 幻灯片 8: AI对话 + RAG ==========
print("创建 Slide 8: AI对话+RAG知识库")
slide = add_blank_slide()
add_slide_title(slide, '核心功能 ③④ AI对话 & RAG知识库', '朋友式AI + 个性化知识增强')
add_page_number(slide, 8)

# 左：AI对话
add_textbox(slide, 1.0, 2.0, 5.5, 0.5, '🤖 小蓝AI — 朋友式对话', font_size=18, color=DARK_BLUE, bold=True)
chat_features = [
    '温暖活泼的校园学习搭子人设',
    '年轻化用语，拉近用户距离',
    '双Tab设计：聊天 / 补课包一键切换',
    'SSE流式输出，逐字实时呈现',
    '支持语音输入(STT) + 语音播报(TTS)',
    'Ctrl+K 快捷键唤起',
]
for i, feat in enumerate(chat_features):
    add_textbox(slide, 1.0, 2.6 + i * 0.45, 5.5, 0.4, f'• {feat}', font_size=13, color=DARK_GRAY)

# 右：RAG知识库
add_accent_bar(slide, 7.0, 2.0, 0.03, 4.5, LIGHT_GRAY)
add_textbox(slide, 7.5, 2.0, 5.5, 0.5, '📚 RAG知识库', font_size=18, color=DARK_BLUE, bold=True)
rag_features = [
    'PDF/DOCX/TXT/MD 文档上传解析',
    '自动分块(~1000字/块, 200字重叠)',
    '向量相似度检索 + 关键词兜底',
    '笔记与知识库自动同步联动',
    '对话时可选择是否启用RAG增强',
    'embedding模型可独立配置',
]
for i, feat in enumerate(rag_features):
    add_textbox(slide, 7.5, 2.6 + i * 0.45, 5.5, 0.4, f'• {feat}', font_size=13, color=DARK_GRAY)

# 底部
add_textbox(slide, 1.0, 6.2, 11.3, 0.7,
            '💡 笔记 ⇄ 知识库双向联动：笔记增删改自动同步RAG，对话回答始终基于用户最新学习资料',
            font_size=13, color=ORANGE, bold=True, alignment=PP_ALIGN.CENTER)

# ========== 幻灯片 9: 补课包 + 学习周报 ==========
print("创建 Slide 9: 补课包+学习周报")
slide = add_blank_slide()
add_slide_title(slide, '核心功能 ⑤⑥ 补课包 & 学习周报', 'AI驱动的个性化学习补救与复盘')
add_page_number(slide, 9)

# 左：补课包
add_card(slide, 1.0, 2.0, 5.5, 4.5,
         '🏫 逃课补课包',
         '素材选择：自由勾选关联笔记与文档\n'
         '智能生成：知识点解析 + 常见误区\n'
         '复习方法：时间分配建议 + 验收标准\n'
         '自测题目：含答案解析，巩固知识点\n'
         '追问功能：生成后可继续向AI提问\n'
         '保存复用：一键保存为笔记长期查看\n\n'
         '📍 不是单次生成，是可迭代的完整学习闭环',
         ORANGE)

# 右：学习周报
add_card(slide, 7.0, 2.0, 5.5, 4.5,
         '📊 学习周报',
         '自动聚合：本周笔记产出量统计\n'
         'DDL追踪：已完成/待完成/已过期\n'
         '学习趋势：知识点覆盖与薄弱分析\n'
         '建议推送：基于数据的学习建议\n'
         '一键导出：生成可分享的学习报告\n\n'
         '📍 从"学了什么"到"该怎么学"的\n'
         '   数据驱动学习决策',
         BLUE)

# ========== 幻灯片 10: 安全与可靠性 ==========
print("创建 Slide 10: 安全与可靠性")
slide = add_blank_slide()
add_slide_title(slide, '安全与可靠性设计', '全方位保障用户体验与数据安全')
add_page_number(slide, 10)

# 5个安全维度
sec_items = [
    ('🔐 身份认证', 'BCrypt密码加密\nJWT双令牌(120min+7天)\n前端401自动刷新', BLUE),
    ('🛡️ 输入安全', 'Bean Validation参数校验\nInputSanitizer清洗\n前端HTML转义防XSS', DARK_BLUE),
    ('🚦 流量控制', '登录/注册: 10次/分钟\nAI/RAG/DDL: 30次/分钟\n防恶意调用', LIGHT_BLUE),
    ('⚠️ 异常处理', '统一异常捕获与响应格式\nApiException/AiServiceException\nRateLimitException分类处理', ORANGE),
    ('🔒 数据隔离', 'DTO脱敏屏蔽内部字段\nuserId/ragDocumentId不暴露\n前端数据最小化原则', BLUE),
]
for i, (title, body, color) in enumerate(sec_items):
    x = 0.6 + i * 2.45
    add_card(slide, x, 2.0, 2.3, 4.5, title, body, color)

# ========== 幻灯片 11: 页面展示 ==========
print("创建 Slide 11: 页面展示")
slide = add_blank_slide()
add_slide_title(slide, '产品界面展示', '从概念到落地，直观呈现核心交互界面')
add_page_number(slide, 11)

# 6个界面占位描述
screens = [
    ('登录与首页', '卡片式数据聚合\n笔记产出 + DDL进度\n过期任务色块预警'),
    ('智能笔记创建', '手写/拍照/上传\nMarkdown实时预览\n文件夹树归档'),
    ('DDL管理面板', 'AI解析批量输入\n智能分类筛选\n已完成/已过期视图'),
    ('小蓝AI对话', '朋友式交互界面\n双Tab切换\n流式Markdown渲染'),
    ('RAG知识库', '文档上传与管理\n检索增强对话\n笔记自动同步'),
    ('学习周报', '数据可视化呈现\n学习趋势分析\n个性化建议推送'),
]
for i, (name, desc) in enumerate(screens):
    x = 0.6 + (i % 3) * 4.1
    y = 2.0 + (i // 3) * 2.6
    add_shape_bg(slide, x, y, 3.8, 2.3, BG_WHITE, MSO_SHAPE.ROUNDED_RECTANGLE)
    add_accent_bar(slide, x, y, 3.8, 0.06, BLUE)
    add_textbox(slide, x + 0.2, y + 0.25, 3.4, 0.4, name, font_size=15,
                color=DARK_BLUE, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x + 0.2, y + 0.8, 3.4, 1.3, desc, font_size=12,
                color=MID_GRAY, alignment=PP_ALIGN.CENTER)
    # 界面占位提示
    add_textbox(slide, x + 0.2, y + 1.05, 3.4, 0.3,
                '（答辩时切到实际页面演示）', font_size=9,
                color=LIGHT_BLUE, alignment=PP_ALIGN.CENTER)

# ========== 幻灯片 12: 未来展望 ==========
print("创建 Slide 12: 未来展望")
slide = add_blank_slide()
add_slide_title(slide, '未来展望', '持续进化，打造更强大的校园AI学习生态')
add_page_number(slide, 12)

roadmap = [
    ('🚀 近期', [
        '迁移至vivo快应用平台，降低使用门槛',
        '引入Redis：分布式限流 + 缓存 + JWT黑名单',
        '完善AI调用全链路日志与可观测性',
    ], BLUE),
    ('🔭 中期', [
        '开放API，赋能更多校园应用接入AI能力',
        'Flyway数据库版本精细化管理',
        '多端适配：手机/平板/PC全覆盖',
    ], LIGHT_BLUE),
    ('🌟 远期', [
        '构建校园学习数据中台，驱动个性化推荐',
        'AI学习教练：自适应学习路径规划',
        '打造开放生态：课程、笔记、题库互联互通',
    ], ORANGE),
]
for i, (phase, items, color) in enumerate(roadmap):
    x = 1.0 + i * 3.8
    # 阶段标题
    add_shape_bg(slide, x, 2.2, 3.5, 0.6, color, MSO_SHAPE.ROUNDED_RECTANGLE)
    add_textbox(slide, x, 2.25, 3.5, 0.5, phase, font_size=18,
                color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    # 内容
    for j, item in enumerate(items):
        add_textbox(slide, x + 0.2, 3.2 + j * 0.9, 3.1, 0.8,
                    f'{j+1}. {item}', font_size=13, color=DARK_GRAY)

# ========== 幻灯片 13: 感谢页 ==========
print("创建 Slide 13: 感谢页")
slide = add_blank_slide()
add_shape_bg(slide, 0, 0, 13.333, 7.5, DARK_BLUE)
add_accent_bar(slide, 3.0, 2.2, 7.333, 0.03, LIGHT_BLUE)
add_accent_bar(slide, 3.0, 5.0, 7.333, 0.03, LIGHT_BLUE)

add_textbox(slide, 0, 2.8, 13.333, 1.5, '感谢聆听',
            font_size=56, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(slide, 0, 4.1, 13.333, 0.8, 'Q & A',
            font_size=28, color=LIGHT_BLUE, alignment=PP_ALIGN.CENTER)
add_textbox(slide, 0, 5.5, 13.333, 0.6,
            '欢迎各位评委老师提出宝贵意见与建议',
            font_size=16, color=WHITE, alignment=PP_ALIGN.CENTER)

# ========== 保存 ==========
output = '蓝心校园AI管家-答辩压缩版.pptx'
prs.save(output)
print(f'\n✅ 完成！文件: {output}')
print(f'   共 {len(prs.slides)} 页幻灯片')
print(f'   配色：保留原Office主题 (蓝色系)')
